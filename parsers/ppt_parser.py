import logging
from pathlib import Path
from typing import Dict, Any

logger = logging.getLogger(__name__)


class PPTParser:
    async def analyze_presentation(self, file_path: str) -> Dict[str, Any]:
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            slides_data = []
            for slide in prs.slides:
                slide_info = {"has_title": slide.has_title,
                              "title": slide.shapes.title.text if slide.has_title else "",
                              "shape_count": len(slide.shapes),
                              "text_count": sum(len(s.text_frame.text) for s in slide.shapes
                                                if hasattr(s, 'text_frame') and s.text_frame.text)}
                slides_data.append(slide_info)
            return {"file_path": file_path, "type": "powerpoint", "slide_count": len(prs.slides),
                    "slides": slides_data,
                    "metadata": {"title": prs.core_properties.title or "",
                                 "author": prs.core_properties.author or "",
                                 "subject": prs.core_properties.subject or ""},
                    "analysis": {"structure_score": min(100, len(prs.slides) * 10),
                                 "has_title_slide": True, "has_agenda": any(
                        "agenda" in (s.shapes.title.text.lower() if s.has_title else "")
                        for s in prs.slides[:3])}}
        except Exception as e:
            logger.warning(f"PPT parse failed: {e}")
            return {"file_path": file_path, "type": "powerpoint", "error": str(e), "slide_count": 0}
